import re
import logging
from io import BytesIO
from pdf2image import convert_from_bytes
from pptx import Presentation
from pdfminer.high_level import extract_text
from PyPDF2 import PdfMerger
from sentence_transformers import SentenceTransformer, util
import streamlit as st

logging.basicConfig(level=logging.INFO)


def extract_text_from_ppt(file):
    prs = Presentation(file)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_text.append({"slide_number": i + 1, "text": " ".join(slide_text)})
    return slides_text

def extract_text_from_pdf(file):
    pdf_text = extract_text(file)
    pages = pdf_text.split("\f")
    pages_text = []
    for i, page in enumerate(pages):
        pages_text.append({"slide_number": i + 1, "text": page})
    return pages_text

def merge_pdfs(uploaded_files):
    merger = PdfMerger()
    merged_file = BytesIO()
    for file in uploaded_files:
        merger.append(file)
    merger.write(merged_file)
    merger.close()
    merged_file.seek(0)
    return merged_file

# Load SentenceTransformer model
@st.cache_resource
def load_model():
    return SentenceTransformer('paraphrase-mpnet-base-v2')
 
 
def preprocess_text(text):
    text = re.sub(r'[\–\-\•]\s*\n', ' ', text)
    text = re.sub(r'^\s*[\•\-\–]\s*', '', text, flags=re.MULTILINE)
    text = text.lower()
    text = re.sub(r'[^\w\s\-\.\,\'\"]', '', text)
    return text


def sliding_window(text, window_size, step_size):
    words = text.split()
    for i in range(0, len(words) - window_size + 1, step_size):
        yield " ".join(words[i:i + window_size])

def get_embeddings(text, model):
    text = preprocess_text(text)
    return model.encode(text, convert_to_tensor=True)

def is_negated(statement):
    negation_words = {"not", "no", "never", "none", "nobody", "nothing", "neither", "nowhere", "hardly", "scarcely", "barely"}
    statement_words = set(statement.lower().split())
    return bool(negation_words & statement_words)

def check_statements(slides, statements, model, window_size=30, step_size=10):
    results = {}
    similarity_threshold_true = 0.6
    similarity_threshold_negation = 0.55

    for statement in statements:
        logging.info(f"Checking statement: {statement}")
        negated = is_negated(statement)
        best_match = {"slide_number": None, "similarity": 0, "window_text": ""}
        statement_embedding = get_embeddings(statement, model)

        for slide in slides:
            slide_number = slide["slide_number"]
            for window_text in sliding_window(slide["text"], window_size, step_size):
                window_embedding = get_embeddings(window_text, model)
                similarity = util.pytorch_cos_sim(window_embedding, statement_embedding)
                similarity_value = similarity.item()
                logging.info(f"Comparing with window text: {window_text}, similarity: {similarity_value}")

                if similarity_value > best_match["similarity"]:
                    best_match = {"slide_number": slide_number, "similarity": similarity_value, "window_text": window_text}

        if negated:
            if best_match["similarity"] >= similarity_threshold_negation:
                results[statement] = {"result": 'False', "similarity": best_match["similarity"], "slide_number": best_match["slide_number"], "window_text": best_match["window_text"]}
            else:
                results[statement] = {"result": 'True', "similarity": best_match["similarity"], "slide_number": best_match["slide_number"], "window_text": best_match["window_text"]}
        else:
            if best_match["similarity"] >= similarity_threshold_true:
                results[statement] = {"result": 'True', "similarity": best_match["similarity"], "slide_number": best_match["slide_number"], "window_text": best_match["window_text"]}
            else:
                results[statement] = {"result": 'False', "similarity": best_match["similarity"], "slide_number": best_match["slide_number"], "window_text": best_match["window_text"]}

    return results

def get_pdf_image(file, page_number):
    try:
        images = convert_from_bytes(file.getvalue(), first_page=page_number, last_page=page_number)
        img_byte_arr = BytesIO()
        images[0].save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)
        return img_byte_arr
    except Exception as e:
        logging.error(f"Error converting PDF to image: {e}")
        return None

def combine_model_results(slides, statements, model, use_gpt=False, lower_threshold=0.40, upper_threshold=0.60):
    results = check_statements(slides, statements, model)
    
    if use_gpt:
        from gpt_utils import check_statement_with_gpt
        for statement in statements:
            similarity_value = results[statement]['similarity']
            if lower_threshold <= similarity_value <= upper_threshold:
                gpt_result = check_statement_with_gpt(statement)
                if 'true' in gpt_result.lower():
                    results[statement]['result'] = 'True'
                elif 'false' in gpt_result.lower():
                    results[statement]['result'] = 'False'
                results[statement]['gpt_result'] = gpt_result
    
    return results

